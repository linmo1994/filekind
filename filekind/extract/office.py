from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from filekind.models import FileRecord, RuntimeConfig


def _office_error(record: FileRecord, method: str, message: str) -> FileRecord:
    record.extract_method = method
    record.pages_extracted = 0
    record.raw_snippet = ""
    record.reason = message
    return record


def extract_docx(record: FileRecord, runtime: RuntimeConfig) -> FileRecord:
    try:
        doc = Document(Path(record.path))
    except Exception:
        return _office_error(record, "docx_error", "无法读取 Word 文件（可能已损坏或格式不符）")

    chunks: list[str] = []
    char_count = 0
    limit = runtime.text_fallback_chars
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        chunks.append(text)
        char_count += len(text)
        if char_count >= limit:
            break
    record.raw_snippet = "\n".join(chunks)[:limit]
    record.pages_extracted = 1
    record.extract_method = "docx"
    return record


def extract_pptx(record: FileRecord, runtime: RuntimeConfig) -> FileRecord:
    try:
        prs = Presentation(Path(record.path))
    except Exception:
        return _office_error(record, "pptx_error", "无法读取 PPT 文件（可能已损坏或格式不符）")

    max_slides = runtime.extract_max_pages
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides):
        if idx >= max_slides:
            break
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append("\n".join(slide_text))
    snippet = "\n\n".join(parts)
    record.raw_snippet = snippet[: runtime.text_fallback_chars]
    record.pages_extracted = min(max_slides, len(prs.slides))
    record.extract_method = "pptx"
    return record


def extract_xlsx(record: FileRecord, runtime: RuntimeConfig) -> FileRecord:
    try:
        wb = load_workbook(Path(record.path), read_only=True, data_only=True)
    except Exception:
        return _office_error(record, "xlsx_error", "无法读取 Excel 文件（可能已损坏或格式不符）")

    try:
        ws = wb.active
        rows: list[str] = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= runtime.excel_max_rows:
                break
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append("\t".join(cells))
    finally:
        wb.close()

    snippet = "\n".join(rows)[: runtime.text_fallback_chars]
    record.raw_snippet = snippet
    record.pages_extracted = 1
    record.extract_method = "xlsx"
    return record
